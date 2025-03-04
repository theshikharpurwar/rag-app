# D:\rag-app\python\compute_embeddings.py

import os
import sys
import json
import logging
import glob
import fitz  # PyMuPDF
from PIL import Image
import io
import argparse
import numpy as np
from qdrant_client import QdrantClient
from qdrant_client.models import VectorParams, Distance, PointStruct
import pptx  # For PowerPoint files

# Configure logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(name)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

def extract_from_pdf(file_path):
    """Extract text and images from a PDF using PyMuPDF."""
    result = []

    try:
        # Open the PDF file
        doc = fitz.open(file_path)
        page_count = len(doc)
        logger.info(f"Processing PDF with {page_count} pages")

        # Process each page
        for page_num in range(page_count):
            page = doc[page_num]

            # Extract text
            text = page.get_text()

            # Create PIL image from page
            pix = page.get_pixmap(alpha=False)
            img_data = pix.tobytes("jpeg")
            img = Image.open(io.BytesIO(img_data))

            result.append({
                'page': page_num + 1,  # 1-based page numbering
                'text': text,
                'image': img
            })

        return result, page_count
    except Exception as e:
        logger.error(f"Error extracting content from PDF: {str(e)}")
        raise

def extract_from_pptx(file_path):
    """Extract text and images from a PowerPoint file."""
    result = []

    try:
        # Open the PowerPoint file
        presentation = pptx.Presentation(file_path)
        slide_count = len(presentation.slides)
        logger.info(f"Processing PowerPoint with {slide_count} slides")

        # Process each slide
        for slide_num, slide in enumerate(presentation.slides):
            # Extract text from all shapes
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texts.append(shape.text)

            text = "\n".join(texts)

            # Create a blank image for the slide (we can't easily render slides as images)
            # In a real solution, you'd use a library that can render PowerPoint slides
            img = Image.new('RGB', (800, 600), color='white')

            result.append({
                'page': slide_num + 1,  # 1-based slide numbering
                'text': text,
                'image': img
            })

        return result, slide_count
    except Exception as e:
        logger.error(f"Error extracting content from PowerPoint: {str(e)}")
        raise

def extract_content(file_path):
    """Extract content from various file types."""
    file_ext = os.path.splitext(file_path)[1].lower()

    if file_ext == '.pdf':
        return extract_from_pdf(file_path)
    elif file_ext in ['.pptx', '.ppt']:
        return extract_from_pptx(file_path)
    else:
        raise ValueError(f"Unsupported file type: {file_ext}")

def create_qdrant_collection(client, collection_name, vector_size):
    """Create a Qdrant collection if it doesn't exist."""
    try:
        collections = client.get_collections().collections
        collection_names = [collection.name for collection in collections]

        if collection_name not in collection_names:
            client.create_collection(
                collection_name=collection_name,
                vectors_config=VectorParams(size=vector_size, distance=Distance.COSINE)
            )
            logger.info(f"Created new collection: {collection_name}")
        else:
            logger.info(f"Collection {collection_name} already exists")
    except Exception as e:
        logger.error(f"Error creating collection: {str(e)}")
        raise

def store_embeddings(embeddings, collection_name):
    """Store embeddings in Qdrant."""
    try:
        # Connect to Qdrant
        client = QdrantClient(host="localhost", port=6333)

        # Get embedding size from first embedding
        if not embeddings:
            raise ValueError("No embeddings to store")

        vector_size = len(embeddings[0]['vector'])
        logger.info(f"Embedding vector size: {vector_size}")

        # Create collection if it doesn't exist
        create_qdrant_collection(client, collection_name, vector_size)

        # Prepare points for insertion
        points = []
        for i, item in enumerate(embeddings):
            # Convert embedding to list if it's a numpy array
            vector = item['vector'].tolist() if isinstance(item['vector'], np.ndarray) else item['vector']

            # Create point with ID based on index
            point = PointStruct(
                id=i,
                vector=vector,
                payload={
                    'text': item['text'],
                    'page': str(item['page'])
                }
            )
            points.append(point)

        # Insert points into collection
        client.upsert(
            collection_name=collection_name,
            points=points
        )

        logger.info(f"Stored {len(embeddings)} embeddings in collection {collection_name}")
        return True
    except Exception as e:
        logger.error(f"Error storing embeddings: {str(e)}")
        raise

def main():
    try:
        # Parse command-line arguments
        parser = argparse.ArgumentParser(description='Compute embeddings for document')
        parser.add_argument('file_path', help='Path to document file (PDF, PPTX)')
        parser.add_argument('--collection_name', default='documents', help='Qdrant collection name')
        parser.add_argument('--model_name', default='clip', help='Embedding model name')
        parser.add_argument('--model_path', default='ViT-B/32', help='Path to embedding model')

        args = parser.parse_args()

        # Log the arguments
        logger.info(f"Processing file: {args.file_path}")
        logger.info(f"Collection name: {args.collection_name}")
        logger.info(f"Model name: {args.model_name}")
        logger.info(f"Model path: {args.model_path}")

        # Extract content from document
        page_contents, page_count = extract_content(args.file_path)

        # Get the embedder
        from embeddings.embed_factory import get_embedder
        embedder = get_embedder(args.model_name, args.model_path)

        # Generate embeddings for each page
        embeddings = []
        for page_content in page_contents:
            # Get embedding for the image
            embedding = embedder.get_embedding(page_content['image'])

            # Store embedding with text and page metadata
            embeddings.append({
                'vector': embedding,
                'text': page_content['text'],
                'page': page_content['page']
            })

        # Store embeddings in Qdrant
        success = store_embeddings(embeddings, args.collection_name)

        # Output result as JSON
        result = {
            'success': success,
            'message': f"Successfully processed {page_count} pages and stored {len(embeddings)} embeddings",
            'pageCount': page_count
        }
        print(json.dumps(result))

    except Exception as e:
        error_message = str(e)
        logger.error(f"Error: {error_message}")

        # Output error as JSON
        error_result = {
            'success': False,
            'message': error_message
        }
        print(json.dumps(error_result))
        sys.exit(1)

if __name__ == "__main__":
    main()